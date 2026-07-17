#!/usr/bin/env python3
"""
build_plano_acao_pptx.py — Gera o Plano de Ação em PPTX a partir do template.

Estratégia: clona `assets/plano_acao_template.pptx` (caso Marcelo) e substitui
o conteúdo textual de cada slide preservando layout, fontes e cores. Isso
garante que a identidade visual da Mentoria de Automações Inteligentes seja
mantida automaticamente.

Uso:
    python build_plano_acao_pptx.py --template assets/plano_acao_template.pptx \\
        --json caminho/plano.json --saida caminho/saida.pptx

Estrutura do JSON em EXEMPLO_PAYLOAD no rodapé deste arquivo.

Dependências: python-pptx
"""

from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Pt


# ----- helpers ---------------------------------------------------------------

def _iter_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame


def _runs_in_order(text_frame) -> list:
    """Retorna todos os runs do text_frame em ordem (linhas → paragrafos → runs)."""
    runs = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            runs.append(run)
    return runs


def _set_paragraph_text_preserving_format(text_frame, lines: list[str]) -> None:
    """
    Substitui linhas de texto preservando o formato do primeiro run de cada
    parágrafo. Quando o número de novas linhas é diferente do número de
    parágrafos existentes, ajusta (sobra → limpa; falta → duplica formato do
    último parágrafo).
    """
    existing = list(text_frame.paragraphs)
    n_existing = len(existing)
    n_new = len(lines)

    for i, line in enumerate(lines):
        if i < n_existing:
            para = existing[i]
            # preservar primeiro run, limpar o resto
            if para.runs:
                first_run = para.runs[0]
                first_run.text = line
                # remover runs extras
                for run in list(para.runs[1:]):
                    run._r.getparent().remove(run._r)
            else:
                para.text = line
        else:
            # adicionar novo parágrafo copiando formato do último existente
            new_para = text_frame.add_paragraph()
            new_para.text = line
            if existing:
                src = existing[-1]
                if src.runs and new_para.runs:
                    src_run = src.runs[0]
                    dst_run = new_para.runs[0]
                    if src_run.font.size:
                        dst_run.font.size = src_run.font.size
                    if src_run.font.name:
                        dst_run.font.name = src_run.font.name
                    if src_run.font.bold is not None:
                        dst_run.font.bold = src_run.font.bold
                    try:
                        if src_run.font.color and src_run.font.color.rgb:
                            dst_run.font.color.rgb = src_run.font.color.rgb
                    except Exception:
                        pass

    # remover parágrafos sobrantes
    if n_new < n_existing:
        for para in existing[n_new:]:
            para._p.getparent().remove(para._p)


def _replace_in_slide(slide, mapping: dict[str, str]) -> None:
    """Substitui ocorrências de chaves do mapping nos text_frames do slide."""
    for tf in _iter_text_frames(slide):
        for para in tf.paragraphs:
            for run in para.runs:
                for k, v in mapping.items():
                    if k in run.text:
                        run.text = run.text.replace(k, v)


# ----- builders por slide ----------------------------------------------------

def _apply_cover(slide, data) -> None:
    """Slide 1 — Capa."""
    nome = data['cliente']
    empresa = data['empresa']
    tags = data.get('tags_prioridade', [])

    for tf in _iter_text_frames(slide):
        text = tf.text
        if 'Marcelo Coelho' in text and 'Coelho Music' not in text:
            _set_paragraph_text_preserving_format(tf, ['Plano de ação', nome])
        elif 'Coelho Music' in text:
            _set_paragraph_text_preserving_format(tf, [empresa, 'Mentoria Automações Inteligentes'])
        elif 'Dashboard pedagógico primeiro' in text or 'Comercial e conteúdo' in text:
            _set_paragraph_text_preserving_format(tf, tags if tags else [text])


def _apply_text_replace(slide, replacements: list[tuple[str, str]]) -> None:
    """Substituições simples mantendo formato exato."""
    mapping = {k: v for k, v in replacements}
    _replace_in_slide(slide, mapping)


def _apply_slide_negocio(slide, data) -> None:
    """Slide 2 — O que entendemos do negócio."""
    bullets = data['entendimento_bullets']  # lista de strings
    stats = data['stats']  # [{numero, label}, ...] com 3 itens

    # Coletar text_frames de bullets (têm marcador • ciano) e os blocos de stats
    bullet_pairs: list[tuple] = []
    stat_blocks: list[tuple] = []

    shapes = list(slide.shapes)
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        # heurística: linhas curtas de número + descrição
        if text in ('41h', '15h', '3') or (len(text) <= 4 and any(c.isdigit() for c in text)):
            stat_blocks.append(('numero', tf))
        elif text.startswith('aprox.') or text.startswith('dedicadas') or text.startswith('frentes'):
            stat_blocks.append(('label', tf))
        elif 'Jazz Corporativo' in text or 'Workshop de criatividade' in text \
                or 'Startup educacional' in text:
            bullet_pairs.append(tf)

    # Substituir bullets: cada bullet ocupa um text_frame separado no template
    for i, tf in enumerate(bullet_pairs):
        if i < len(bullets):
            _set_paragraph_text_preserving_format(tf, [bullets[i]])

    # Substituir stats (3 pares numero+label, em ordem visual)
    numero_tfs = [tf for k, tf in stat_blocks if k == 'numero']
    label_tfs = [tf for k, tf in stat_blocks if k == 'label']
    for i, st in enumerate(stats[:3]):
        if i < len(numero_tfs):
            _set_paragraph_text_preserving_format(numero_tfs[i], [str(st['numero'])])
        if i < len(label_tfs):
            _set_paragraph_text_preserving_format(label_tfs[i], [st['label']])


def _replace_bullets_by_marker(slide, marker_substrings: Iterable[str],
                                new_bullets: list[str]) -> None:
    """Substitui text_frames de bullets que começam com determinadas substrings."""
    matched_tfs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        for sub in marker_substrings:
            if sub in text:
                matched_tfs.append(tf)
                break
    for i, tf in enumerate(matched_tfs):
        if i < len(new_bullets):
            _set_paragraph_text_preserving_format(tf, [new_bullets[i]])


def _apply_aula_link(slide, hyperlinks: list[dict]) -> None:
    """Adiciona hyperlinks nos itens 'abrir' verdes do slide de aulas."""
    abrir_tfs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if tf.text.strip().lower() == 'abrir':
            abrir_tfs.append(tf)

    for i, link_info in enumerate(hyperlinks):
        if i < len(abrir_tfs) and link_info.get('url'):
            tf = abrir_tfs[i]
            for para in tf.paragraphs:
                for run in para.runs:
                    rPr = run._r.get_or_add_rPr()
                    # configurar hyperlink
                    from pptx.oxml.ns import qn
                    hl = rPr.find(qn('a:hlinkClick'))
                    if hl is None:
                        from lxml import etree
                        hl = etree.SubElement(rPr, qn('a:hlinkClick'))
                    # adicionar relationship
                    part = slide.part
                    rId = part.relate_to(
                        link_info['url'],
                        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
                        is_external=True,
                    )
                    hl.set(qn('r:id'), rId)


# ----- builder principal -----------------------------------------------------

def build(template_path: Path, data: dict, saida: Path) -> None:
    prs = Presentation(str(template_path))
    slides = list(prs.slides)
    if len(slides) < 12:
        print(f'AVISO: template tem apenas {len(slides)} slides — esperava 12. Procedendo mesmo assim.')

    # Slide 1 — Capa
    _apply_cover(slides[0], data)

    # Slide 2 — Entendimento do negócio
    if len(slides) > 1:
        _apply_slide_negocio(slides[1], data)

    # Slide 3 — Gargalo central
    if len(slides) > 2 and data.get('gargalo'):
        s = slides[2]
        # frase âncora
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'arquitetura de processos' in tf.text.lower():
                _set_paragraph_text_preserving_format(tf, [data['gargalo']['frase_ancora']])
        _replace_bullets_by_marker(
            s,
            ['Construir e amadurecer', 'Organizar indicadores',
             'Automatizar tráfego', 'Padronizar o uso de IA'],
            data['gargalo']['bullets'],
        )

    # Slide 4 — Por que a #1 vem primeiro
    if len(slides) > 3 and data.get('por_que_primeiro'):
        s = slides[3]
        # 4 cards: mini-título + descrição. Vou substituir pelas posições conhecidas
        cards = data['por_que_primeiro']['cards']  # [{titulo, descricao}, ...]
        # mini-títulos
        titulo_subs = ['Já está em andamento', 'Maior carga semanal',
                       'Prova de valor', 'Base comercial']
        # encontrar tfs por substring
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            txt = tf.text.strip()
            for i, t in enumerate(titulo_subs):
                if t == txt and i < len(cards):
                    _set_paragraph_text_preserving_format(tf, [cards[i]['titulo']])

        # descrições — heurística: cinza Arial pequeno
        desc_subs = ['Aproveita energia atual', 'É a frente com mais tempo',
                     'Relatórios podem demonstrar', 'Indicadores fortalecem proposta']
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            txt = tf.text.strip()
            for i, t in enumerate(desc_subs):
                if t in txt and i < len(cards):
                    _set_paragraph_text_preserving_format(tf, [cards[i]['descricao']])

        # frase final
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'Decisão de prioridade:' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [f"Decisão de prioridade: {data['por_que_primeiro']['decisao']}"])

    # Slide 5 — Estrutura recomendada
    if len(slides) > 4 and data.get('estrutura_projeto'):
        s = slides[4]
        # título do slide
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'Estrutura recomendada' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [f"Estrutura recomendada do {data['estrutura_projeto'].get('projeto', 'projeto')}"])
                break
        _replace_bullets_by_marker(
            s,
            ['Partir de perguntas', 'Separar indicadores',
             'Conectar cada indicador', 'Começar simples'],
            data['estrutura_projeto']['bullets'],
        )
        # pergunta-chave em verde
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'Pergunta-chave' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [f"Pergunta-chave: {data['estrutura_projeto']['pergunta_chave']}"])

    # Slide 6 — Roadmap
    if len(slides) > 5 and data.get('roadmap'):
        s = slides[5]
        prazos = data['roadmap']  # [{prazo, descricao}, ...]
        prazo_subs = ['7 dias', '30 dias', '60 dias', '90 dias']
        desc_subs = ['Diagnóstico no portal', 'KPIs, fontes de dados',
                     'Modelo de relatório', 'Automação de relatórios']
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            txt = tf.text.strip()
            for i, t in enumerate(prazo_subs):
                if t == txt and i < len(prazos):
                    _set_paragraph_text_preserving_format(tf, [prazos[i]['prazo']])
            for i, t in enumerate(desc_subs):
                if t in txt and i < len(prazos):
                    _set_paragraph_text_preserving_format(tf, [prazos[i]['descricao']])

    # Slide 7 — Aulas começar agora (8 itens 01-08)
    if len(slides) > 6 and data.get('aulas_inicio'):
        s = slides[6]
        # substituir nomes/desc de cada aula
        aulas = data['aulas_inicio']
        # vamos pegar pares (nome, descricao) na ordem visual
        # heurística: text frames bold branco = nome, text frames cinza = descricao
        nome_tfs = []
        desc_tfs = []
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if not text or text == 'abrir':
                continue
            if text == 'Aulas para começar agora':
                continue
            if text.isdigit() or (len(text) == 2 and text.startswith('0')):
                continue
            if 'PLANO DE AÇÃO' in text:
                continue
            # heurística simples por cor do primeiro run
            try:
                first = tf.paragraphs[0].runs[0]
                if first.font.bold:
                    nome_tfs.append(tf)
                else:
                    desc_tfs.append(tf)
            except (IndexError, AttributeError):
                pass

        for i, aula in enumerate(aulas[:8]):
            if i < len(nome_tfs):
                _set_paragraph_text_preserving_format(nome_tfs[i], [aula['nome']])
            if i < len(desc_tfs):
                _set_paragraph_text_preserving_format(desc_tfs[i], [aula['descricao']])

        # hyperlinks
        _apply_aula_link(s, [{'url': a.get('link', '')} for a in aulas[:8]])

    # Slide 8 — Aulas segunda etapa
    if len(slides) > 7 and data.get('aulas_segunda'):
        s = slides[7]
        aulas = data['aulas_segunda']
        nome_tfs = []
        desc_tfs = []
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if not text or text == 'abrir' or text in ('09', '10', '11', '12'):
                continue
            if 'Sequência recomendada' in text or 'evitar dispersão' in text:
                continue
            if 'PLANO DE AÇÃO' in text or text == 'Aulas para segunda etapa':
                continue
            try:
                first = tf.paragraphs[0].runs[0]
                if first.font.bold:
                    nome_tfs.append(tf)
                else:
                    desc_tfs.append(tf)
            except (IndexError, AttributeError):
                pass

        for i, aula in enumerate(aulas[:4]):
            if i < len(nome_tfs):
                _set_paragraph_text_preserving_format(nome_tfs[i], [aula['nome']])
            if i < len(desc_tfs):
                _set_paragraph_text_preserving_format(desc_tfs[i], [aula['descricao']])
        _apply_aula_link(s, [{'url': a.get('link', '')} for a in aulas[:4]])

    # Slide 9 — Recursos
    if len(slides) > 8 and data.get('recursos'):
        s = slides[8]
        _replace_bullets_by_marker(
            s,
            ['Gestor de Métricas', 'UI/UX Pro Max', 'Notion MCP',
             'Supabase MCP', 'Instagram OS'],
            data['recursos'],
        )

    # Slide 10 — Recorrências
    if len(slides) > 9 and data.get('recorrencias'):
        s = slides[9]
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'aceleração prática' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [data['recorrencias'].get('frase_ancora',
                                                  'Usar as recorrências como ambiente de aceleração prática.')])
        _replace_bullets_by_marker(
            s,
            ['Levar dúvidas', 'Validar indicadores', 'Pedir orientação',
             'Destravar documentação', 'Alinhar próximos passos'],
            data['recorrencias']['bullets'],
        )

    # Slide 11 — Próximos passos
    if len(slides) > 10 and data.get('proximos_passos'):
        s = slides[10]
        passos = data['proximos_passos']['cards']  # [{titulo, descricao}]
        titulo_subs = ['1. Acessar', '2. Listar dados', '3. Validar', '4. Implementar']
        desc_subs = ['Entrar no portal', 'Mapear o que já existe',
                     'Definir indicadores mínimos', 'Levar dúvidas para as recorrências']
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            txt = tf.text.strip()
            for i, t in enumerate(titulo_subs):
                if txt.startswith(t.split(' ')[0]) and i < len(passos):
                    _set_paragraph_text_preserving_format(tf, [passos[i]['titulo']])
            for i, t in enumerate(desc_subs):
                if t in txt and i < len(passos):
                    _set_paragraph_text_preserving_format(tf, [passos[i]['descricao']])
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'Foco inicial:' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [f"Foco inicial: {data['proximos_passos']['foco']}"])

    # Slide 12 — Encerramento
    if len(slides) > 11 and data.get('encerramento'):
        s = slides[11]
        linhas = data['encerramento']['frases']  # 3 strings
        # encontrar 3 text frames com tamanho 36pt branco
        big_tfs = []
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            try:
                first = tf.paragraphs[0].runs[0]
                if first.font.size and first.font.size.pt >= 30:
                    big_tfs.append(tf)
            except (IndexError, AttributeError):
                pass
        for i, line in enumerate(linhas[:3]):
            if i < len(big_tfs):
                _set_paragraph_text_preserving_format(big_tfs[i], [line])
        # mini-texto cinza
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if 'primeiro ciclo da mentoria' in tf.text or 'comercial ganha base' in tf.text:
                _set_paragraph_text_preserving_format(
                    tf, [data['encerramento'].get('explicacao', tf.text)])

    saida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(saida))


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument('--template', required=True, type=Path)
    parser.add_argument('--json', required=True, type=Path)
    parser.add_argument('--saida', required=True, type=Path)
    args = parser.parse_args(argv)

    data = json.loads(args.json.read_text(encoding='utf-8'))
    build(args.template, data, args.saida)
    print(f'OK: {args.saida}')
    return 0


# Estrutura esperada do JSON (exemplo do caso Marcelo)
EXEMPLO_PAYLOAD = {
    'cliente': 'Marcelo Coelho',
    'empresa': 'Coelho Music Produções LTDA',
    'tags_prioridade': [
        'Dashboard pedagógico primeiro.',
        'Comercial e conteúdo em sequência.',
    ],
    'entendimento_bullets': [
        'Jazz Corporativo: palestra sobre criatividade, cultura, colaboração e improviso.',
        'Workshop de criatividade: experiência formativa para empresas, escolas e lideranças.',
        'Startup educacional: formação integral e socioemocional por meio da música.',
    ],
    'stats': [
        {'numero': '41h', 'label': 'aprox. mapeadas em rotinas semanais com potencial de padronização'},
        {'numero': '15h', 'label': 'dedicadas à construção do dashboard'},
        {'numero': '3', 'label': 'frentes de negócio que precisam conversar entre si'},
    ],
    'gargalo': {
        'frase_ancora': 'O desafio não é falta de produto. É arquitetura de processos.',
        'bullets': [
            'Construir e amadurecer o dashboard da startup educacional.',
            'Organizar indicadores, relatórios e evidências de impacto para escolas.',
            'Automatizar tráfego, marketing e comercial para palestra e workshop.',
            'Padronizar o uso de IA para conteúdo, proposta, atendimento, relatório e gestão.',
        ],
    },
    'por_que_primeiro': {
        'cards': [
            {'titulo': 'Já está em andamento', 'descricao': 'Aproveita energia atual e reduz fricção de implementação.'},
            {'titulo': 'Maior carga semanal', 'descricao': 'É a frente com mais tempo mapeado no diagnóstico.'},
            {'titulo': 'Prova de valor', 'descricao': 'Relatórios podem demonstrar impacto para escolas.'},
            {'titulo': 'Base comercial', 'descricao': 'Indicadores fortalecem proposta, conteúdo e narrativa.'},
        ],
        'decisao': 'começar pelo dashboard pedagógico da startup.',
    },
    'estrutura_projeto': {
        'projeto': 'dashboard',
        'bullets': [
            'Partir de perguntas estratégicas, não de gráficos soltos.',
            'Separar indicadores de engajamento, desenvolvimento, operação, satisfação, comercial e impacto.',
            'Conectar cada indicador a uma decisão ou a um trecho do relatório para escolas.',
            'Começar simples: base de dados, KPIs mínimos e primeiras visualizações.',
        ],
        'pergunta_chave': 'que evidência ajuda a escola a perceber valor e continuar avançando?',
    },
    'roadmap': [
        {'prazo': '7 dias', 'descricao': 'Diagnóstico no portal, primeira aula, recorrências e levantamento de dados.'},
        {'prazo': '30 dias', 'descricao': 'KPIs, fontes de dados, estrutura base e MVP do dashboard.'},
        {'prazo': '60 dias', 'descricao': 'Modelo de relatório para escolas e pipeline comercial básico.'},
        {'prazo': '90 dias', 'descricao': 'Automação de relatórios, rotinas de conteúdo e primeiros agentes.'},
    ],
    'aulas_inicio': [],   # lista de {nome, descricao, link}
    'aulas_segunda': [],
    'recursos': [],
    'recorrencias': {
        'frase_ancora': 'Usar as recorrências como ambiente de aceleração prática.',
        'bullets': [],
    },
    'proximos_passos': {
        'cards': [],
        'foco': 'dashboard pedagógico, relatório para escolas e organização dos dados.',
    },
    'encerramento': {
        'frases': ['Agora é transformar', 'o dashboard em prova', 'de valor.'],
        'explicacao': 'O primeiro ciclo da mentoria deve gerar clareza de dados, indicadores e relatório. Depois, o comercial ganha base para crescer.',
    },
}


if __name__ == '__main__':
    sys.exit(main(sys.argv[1:]))
